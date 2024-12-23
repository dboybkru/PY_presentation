from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
from io import BytesIO
from datetime import datetime

# Создание объекта презентации
presentation = Presentation()

# Слайд 1: Заголовок
slide1 = presentation.slides.add_slide(presentation.slide_layouts[0])  # Заголовок
title1 = slide1.shapes.title
subtitle1 = slide1.placeholders[1]
title1.text = "Наша подписка в работе и жизни"
subtitle1.text = "Почему нам стоит ее использовать?\nД. В. Г.\n" + datetime.now().strftime("%d %B %Y")

# Слайд 2: Введение
slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])  # Заголовок и содержание
title2 = slide2.shapes.title
content2 = slide2.placeholders[1]
title2.text = "Введение"
content2.text = (
    "Наша подписка — это сервис, который помогает организовать и управлять подписками на различные ресурсы.\n"
    "Он позволяет пользователям легко отслеживать свои подписки и получать уведомления о новых обновлениях.\n"
    "Важно понять, как наша подписка может улучшить вашу жизнь и работу."
)

# Слайд 3: Преимущества использования нашей подписки
slide3 = presentation.slides.add_slide(presentation.slide_layouts[1])
title3 = slide3.shapes.title
content3 = slide3.placeholders[1]
title3.text = "Преимущества использования нашей подписки"
content3.text = (
    "1. Удобство: все подписки в одном месте.\n"
    "2. Экономия времени: автоматические уведомления о новых материалах.\n"
    "3. Персонализация: настройка интересов для получения релевантного контента."
)

# Слайд 4: Примеры применения нашей подписки
slide4 = presentation.slides.add_slide(presentation.slide_layouts[5])  # Заголовок только
title4 = slide4.shapes.title
title4.text = "Примеры применения нашей подписки"

# Добавление таблицы
rows = 5
cols = 3
table = slide4.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(8), Inches(2)).table

# Заголовки таблицы
table.cell(0, 0).text = "Сфера"
table.cell(0, 1).text = "Применение"
table.cell(0, 2).text = "Преимущества"

# Заполнение таблицы
data = [
    ["Образование", "Подписка на курсы", "Доступ к актуальным материалам"],
    ["Развлечения", "Подписка на стриминговые сервисы", "Неограниченный доступ к контенту"],
    ["Новости", "Подписка на новостные рассылки", "Актуальная информация в реальном времени"],
    ["Спорт", "Подписка на спортивные события", "Уведомления о матчах и результатах"],
]
for i, row in enumerate(data, start=1):
    for j, value in enumerate(row):
        table.cell(i, j).text = value

# Слайд 5: Диаграмма преимуществ
slide5 = presentation.slides.add_slide(presentation.slide_layouts[5])  # Заголовок только
title5 = slide5.shapes.title
title5.text = "Преимущества использования нашей подписки"

# Данные для диаграммы
labels = ['Удобство', 'Экономия времени', 'Персонализация']
sizes = [50, 30, 20]
colors = ['#ff9999','#66b3ff','#99ff99']
explode = (0.1, 0, 0)  # Выделение первого сегмента

# Создание объёмной диаграммы
fig, ax = plt.subplots()
ax.pie(sizes, labels=labels, colors=colors, autopct='%1.1f%%', startangle=90, explode=explode, shadow=True)
ax.axis('equal')  # Равные оси для круговой диаграммы

# Сохранение диаграммы в буфер
image_stream = BytesIO()
plt.savefig(image_stream, format='png')
plt.close(fig)
image_stream.seek(0)

# Вставка диаграммы в слайд
left = Inches(1)
top = Inches(2)
pic = slide5.shapes.add_picture(image_stream, left, top, width=Inches(6), height=Inches(4))

# Слайд 6: Потенциальные риски и этические аспекты
slide6 = presentation.slides.add_slide(presentation.slide_layouts[1])
title6 = slide6.shapes.title
content6 = slide6.placeholders[1]
title6.text = "Потенциальные риски и этические аспекты"
content6.text = (
    "1. Защита данных: необходимость защиты личной информации пользователей.\n"
    "2. Ошибки в алгоритмах: возможные недоразумения в рекомендациях.\n"
    "3. Зависимость от технологий: риск чрезмерного использования сервиса."
)

# Слайд 7: Заключение
slide7 = presentation.slides.add_slide(presentation.slide_layouts[1])
title7 = slide7.shapes.title
content7 = slide7.placeholders[1]
title7.text = "Заключение"
content7.text = (
    "Наша подписка открывает новые горизонты для работы и жизни.\n"
    "Важно адаптироваться к этим изменениям и использовать возможности, которые она предлагает.\n"
    "Призыв к действию: Начните использовать нашу подписку в своей жизни и работе!"
)

# Сохранение презентации
presentation.save("Наша_подписка_в_работе_и_жизни.pptx")